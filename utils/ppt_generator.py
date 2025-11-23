"""
PPT生成模块 - 将图片生成PPT文件
"""
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import math


class PPTGenerator:
    """PPT生成器"""
    
    def __init__(self):
        # PPT标准尺寸（16:9比例）
        # 使用标准16:9尺寸：宽度13.333英寸，高度7.5英寸
        self.slide_width = Inches(13.333)  # 宽度13.333英寸（16:9比例）
        self.slide_height = Inches(7.5)   # 高度7.5英寸
        
        # 左上角边距（美观间距）
        self.margin_left = Inches(0.5)    # 左边距0.5英寸
        self.margin_top = Inches(0.5)      # 上边距0.5英寸
        
        # 图片目标面积（PPT面积的比例）
        # 调整为35%，平衡图片大小和空白区域
        self.target_area_ratio = 0.35  # 35%的面积
    
    def calculate_image_size(self, image_path):
        """
        计算图片在PPT中的尺寸（保持宽高比，面积占PPT的一半）
        
        Args:
            image_path: 图片路径
        
        Returns:
            tuple: (width, height) 单位：Inches
        """
        # 打开图片获取原始尺寸（像素）
        img = Image.open(image_path)
        img_width_px, img_height_px = img.size
        
        # 获取图片DPI（如果有的话），否则使用默认200 DPI
        dpi = img.info.get('dpi', (200, 200))[0] if 'dpi' in img.info else 200
        
        # 将像素转换为英寸
        img_width_inch = img_width_px / dpi
        img_height_inch = img_height_px / dpi
        
        # 计算PPT可用区域（考虑左上角边距）
        # 图片放在左上角，所以只需要考虑左边和上边的边距
        available_width = (self.slide_width / Inches(1)) - (self.margin_left / Inches(1))
        available_height = (self.slide_height / Inches(1)) - (self.margin_top / Inches(1))
        
        # 计算PPT总面积
        slide_area = (self.slide_width / Inches(1)) * (self.slide_height / Inches(1))
        
        # 目标图片面积（PPT面积的一半）
        target_area = slide_area * self.target_area_ratio
        
        # 计算图片原始面积
        img_area = img_width_inch * img_height_inch
        
        if img_area <= 0:
            # 如果图片面积为0，使用默认尺寸
            return Inches(5), Inches(3)
        
        # 计算缩放比例，使图片面积等于目标面积
        # 面积比 = (scale * width) * (scale * height) / (width * height) = scale^2
        # scale^2 = target_area / img_area
        scale = math.sqrt(target_area / img_area)
        
        # 计算缩放后的尺寸
        scaled_width = img_width_inch * scale
        scaled_height = img_height_inch * scale
        
        # 确保图片不超过可用区域（考虑边距）
        max_width = available_width
        max_height = available_height
        
        # 如果缩放后的尺寸超过可用区域，需要进一步缩小
        if scaled_width > max_width or scaled_height > max_height:
            scale_x = max_width / img_width_inch if img_width_inch > 0 else 1.0
            scale_y = max_height / img_height_inch if img_height_inch > 0 else 1.0
            scale = min(scale_x, scale_y)
            scaled_width = img_width_inch * scale
            scaled_height = img_height_inch * scale
        
        return Inches(scaled_width), Inches(scaled_height)
    
    def merge_images(self, image_paths, output_path):
        """
        将多张图片垂直拼接成一张图片
        
        Args:
            image_paths: 图片路径列表
            output_path: 输出图片路径
        """
        if not image_paths:
            raise ValueError("图片列表为空")
        
        if len(image_paths) == 1:
            # 如果只有一张图片，直接复制
            img = Image.open(image_paths[0])
            img.save(output_path, 'PNG')
            return
        
        images = []
        max_width = 0
        total_height = 0
        
        # 打开所有图片并计算尺寸
        for img_path in image_paths:
            img = Image.open(img_path)
            images.append(img)
            max_width = max(max_width, img.width)
            total_height += img.height
        
        # 创建新图片（垂直拼接）
        merged_img = Image.new('RGB', (max_width, total_height), color='white')
        
        # 拼接图片
        y_offset = 0
        for img in images:
            # 居中放置（如果宽度不同）
            x_offset = (max_width - img.width) // 2
            merged_img.paste(img, (x_offset, y_offset))
            y_offset += img.height
        
        # 保存拼接后的图片
        merged_img.save(output_path, 'PNG')
    
    def add_slide_with_image(self, prs, image_path):
        """
        添加包含图片的幻灯片
        
        Args:
            prs: Presentation对象
            image_path: 图片路径
        """
        # 创建新幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 计算图片尺寸
        img_width, img_height = self.calculate_image_size(image_path)
        
        # 图片位置：左上角，但保留适当边距
        left = self.margin_left
        top = self.margin_top
        
        slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)
    
    def create_ppt(self, image_paths, output_path):
        """
        创建PPT文件
        
        Args:
            image_paths: 图片路径列表
            output_path: 输出PPT文件路径
        """
        try:
            # 创建演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸
            prs.slide_width = self.slide_width
            prs.slide_height = self.slide_height
            
            # 为每张图片添加一页幻灯片
            for image_path in image_paths:
                self.add_slide_with_image(prs, image_path)
            
            # 保存PPT
            prs.save(output_path)
        except Exception as e:
            raise Exception(f"PPT生成失败: {str(e)}")

